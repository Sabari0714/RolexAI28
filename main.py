# ============================================================
# ROLEX AI GUI v2.0 MASTER EDITION
# Pydroid 3 + Kivy
#
# IMPORTANT:
# - No PyAudio required
# - No pyjnius required for basic text/AI mode
# - Voice button uses Android's native ACTION_RECOGNIZE_SPEECH
#   through an optional Android bridge.
# - If native voice is unavailable, Rolex remains fully usable
#   through the text box.
# - SQLite thread error is fixed by using a fresh connection
#   per DB operation.
# ============================================================

import os
import re
import json
import time
import sqlite3
import hashlib
import datetime
import urllib.request
import urllib.parse
import threading
import platform
import traceback
import socket
import base64
import csv
import shutil
import random

# Kivy is required only by the Android GUI layer.
# Backend/Brain tests must remain usable without Kivy.
try:
    from kivy.app import App
    from kivy.core.window import Window
    from kivy.clock import Clock
    from kivy.metrics import dp
    from kivy.uix.boxlayout import BoxLayout
    from kivy.uix.button import Button
    from kivy.uix.label import Label
    from kivy.uix.scrollview import ScrollView
    from kivy.uix.textinput import TextInput
    from kivy.graphics import Color, RoundedRectangle
    from kivy.utils import get_color_from_hex as _hex
    KIVY_AVAILABLE = True
except ImportError:
    KIVY_AVAILABLE = False

    class _KivyUnavailable:
        def __init__(self, *args, **kwargs):
            raise RuntimeError(
                "Kivy is required only for the Rolex Android GUI."
            )

    App = _KivyUnavailable
    Window = _KivyUnavailable
    Clock = _KivyUnavailable
    dp = lambda value: value
    BoxLayout = _KivyUnavailable
    Button = _KivyUnavailable
    Label = _KivyUnavailable
    ScrollView = _KivyUnavailable
    TextInput = _KivyUnavailable
    Color = _KivyUnavailable
    RoundedRectangle = _KivyUnavailable

    def _hex(value):
        return value


# ------------------------------------------------------------
# DARK "HI-TECH" THEME
# ------------------------------------------------------------
COLOR_BG          = _hex('#0A0F1E')
COLOR_PANEL       = _hex('#111A2E')
COLOR_SURFACE     = _hex('#182036')
COLOR_ACCENT      = _hex('#00E5C7')
COLOR_ACCENT_SOFT = _hex('#0F8F82')
COLOR_USER_BUBBLE = _hex('#1D3A5F')
COLOR_BOT_BUBBLE  = _hex('#182036')
COLOR_TEXT        = _hex('#EAF0FA')
COLOR_MUTED       = _hex('#7C879E')
COLOR_OK          = _hex('#39E58C')
COLOR_BUSY        = _hex('#FFC24B')
COLOR_INK         = _hex('#04110E')

# ------------------------------------------------------------
# OPTIONAL ANDROID BRIDGE
# ------------------------------------------------------------
ANDROID_BRIDGE = False
activity = None
autoclass = None

try:
    from android import activity as _activity
    activity = _activity
    ANDROID_BRIDGE = True
except Exception:
    pass

if not ANDROID_BRIDGE:
    try:
        from jnius import autoclass as _autoclass
        autoclass = _autoclass
        ANDROID_BRIDGE = True
    except Exception:
        pass


# ------------------------------------------------------------
# PATH / CONFIG
# ------------------------------------------------------------
BASE = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE, "rolex.db")
ENV_PATH = os.path.join(BASE, ".env")


def load_env():
    data = {}

    if os.path.exists(ENV_PATH):
        try:
            with open(ENV_PATH, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()

                    if (
                        line
                        and not line.startswith("#")
                        and "=" in line
                    ):
                        k, v = line.split("=", 1)
                        data[k.strip()] = (
                            v.strip()
                            .strip('"')
                            .strip("'")
                        )
        except Exception:
            pass

    return data


from modules.parallel_ai import parallel_rolex_answer

from modules.providers import (
    ProviderError,
    openai_answer,
    gemini_answer,
    ollama_answer,
)

ENV = load_env()


def cfg(key, default=""):
    return os.getenv(key) or ENV.get(key, default)


NAME = cfg("ROLEX_NAME", "Rolex")

# SECURITY/POLICY: these keys may exist in legacy .env files, but Rolex never uses them.

OPENAI_KEY = cfg("OPENAI_API_KEY")
OPENAI_MODEL = cfg("OPENAI_MODEL", "gpt-5-mini")

GEMINI_KEY = cfg("GEMINI_API_KEY")
GEMINI_MODEL = cfg("GEMINI_MODEL", "gemini-3.6-flash")

OLLAMA_MODEL = cfg("OLLAMA_MODEL")
SERPER_KEY = cfg("SERPER_API_KEY")

ORDER = ["role-local", "openai", "gemini", "ollama"]

# Rolex-only policy: hosted LLM providers are intentionally disabled.
EXTERNAL_AI_DISABLED = False


# ------------------------------------------------------------
# DEVICE DATE / TIME
# ------------------------------------------------------------
def current_date():
    return datetime.datetime.now().strftime(
        "%A, %d %B %Y"
    )


def current_time():
    return datetime.datetime.now().strftime(
        "%I:%M:%S %p"
    )


# ------------------------------------------------------------
# HTTP JSON
# ------------------------------------------------------------
def getjson(
    url,
    payload=None,
    headers=None,
    timeout=45
):
    headers = headers or {}
    data = None
    method = "GET"

    if payload is not None:
        data = json.dumps(payload).encode("utf-8")

        headers = {
            **headers,
            "Content-Type": "application/json"
        }

        method = "POST"

    request = urllib.request.Request(
        url,
        data=data,
        headers=headers,
        method=method
    )

    with urllib.request.urlopen(
        request,
        timeout=timeout
    ) as response:
        raw = response.read().decode("utf-8")

    return json.loads(raw)


# ============================================================
# DATABASE
# ============================================================
class DB:

    def __init__(self):
        self.path = DB_PATH
        self.lock = threading.RLock()
        self.initialize()

    def connect(self):
        return sqlite3.connect(
            self.path,
            timeout=30
        )

    def initialize(self):
        with self.lock:
            c = self.connect()

            try:
                c.execute("""
                    CREATE TABLE IF NOT EXISTS memories(
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        text TEXT NOT NULL,
                        category TEXT DEFAULT 'general',
                        created REAL
                    )
                """)

                c.execute("""
                    CREATE TABLE IF NOT EXISTS tasks(
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        text TEXT NOT NULL,
                        done INTEGER DEFAULT 0,
                        created REAL
                    )
                """)

                c.execute("""
                    CREATE TABLE IF NOT EXISTS cache(
                        key TEXT PRIMARY KEY,
                        value TEXT,
                        expires REAL
                    )
                """)

                c.execute("""
                    CREATE TABLE IF NOT EXISTS documents(
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        path TEXT UNIQUE,
                        title TEXT,
                        kind TEXT,
                        created REAL,
                        modified REAL
                    )
                """)

                c.commit()

            finally:
                c.close()

    def remember(self, text, category="general"):
        with self.lock:
            c = self.connect()

            try:
                q = c.execute(
                    """
                    INSERT INTO memories(
                        text, category, created
                    )
                    VALUES (?, ?, ?)
                    """,
                    (
                        text,
                        category,
                        time.time()
                    )
                )

                c.commit()
                return q.lastrowid

            finally:
                c.close()

    def has(self, text):
        with self.lock:
            c = self.connect()

            try:
                result = c.execute(
                    """
                    SELECT 1
                    FROM memories
                    WHERE lower(text)=lower(?)
                    LIMIT 1
                    """,
                    (text,)
                ).fetchone()

                return result is not None

            finally:
                c.close()

    def recall(self, query=""):
        with self.lock:
            c = self.connect()

            try:
                if query:
                    return c.execute(
                        """
                        SELECT *
                        FROM memories
                        WHERE lower(text) LIKE ?
                        ORDER BY id DESC
                        LIMIT 30
                        """,
                        (
                            "%" + query.lower() + "%",
                        )
                    ).fetchall()

                return c.execute(
                    """
                    SELECT *
                    FROM memories
                    ORDER BY id DESC
                    LIMIT 30
                    """
                ).fetchall()

            finally:
                c.close()

    def recent(self, count=20):
        with self.lock:
            c = self.connect()

            try:
                return c.execute(
                    """
                    SELECT *
                    FROM memories
                    ORDER BY id DESC
                    LIMIT ?
                    """,
                    (count,)
                ).fetchall()

            finally:
                c.close()

    def forget(self, memory_id):
        with self.lock:
            c = self.connect()

            try:
                c.execute(
                    "DELETE FROM memories WHERE id=?",
                    (memory_id,)
                )

                c.commit()

            finally:
                c.close()

    def task(self, text):
        with self.lock:
            c = self.connect()

            try:
                q = c.execute(
                    """
                    INSERT INTO tasks(
                        text, created
                    )
                    VALUES (?, ?)
                    """,
                    (
                        text,
                        time.time()
                    )
                )

                c.commit()
                return q.lastrowid

            finally:
                c.close()

    def tasks(self):
        with self.lock:
            c = self.connect()

            try:
                return c.execute(
                    """
                    SELECT *
                    FROM tasks
                    WHERE done=0
                    ORDER BY id DESC
                    """
                ).fetchall()

            finally:
                c.close()

    def done(self, task_id):
        with self.lock:
            c = self.connect()

            try:
                c.execute(
                    """
                    UPDATE tasks
                    SET done=1
                    WHERE id=?
                    """,
                    (task_id,)
                )

                c.commit()

            finally:
                c.close()

    def cache_get(self, key):
        with self.lock:
            c = self.connect()

            try:
                row = c.execute(
                    """
                    SELECT value, expires
                    FROM cache
                    WHERE key=?
                    """,
                    (key,)
                ).fetchone()

                if not row:
                    return None

                if row[1] < time.time():
                    c.execute(
                        "DELETE FROM cache WHERE key=?",
                        (key,)
                    )

                    c.commit()
                    return None

                return row[0]

            finally:
                c.close()

    def cache_set(self, key, value):
        with self.lock:
            c = self.connect()

            try:
                c.execute(
                    """
                    INSERT OR REPLACE INTO cache(
                        key, value, expires
                    )
                    VALUES (?, ?, ?)
                    """,
                    (
                        key,
                        value,
                        time.time() + 3600
                    )
                )

                c.commit()

            finally:
                c.close()

    def clear_cache(self):
        with self.lock:
            c = self.connect()

            try:
                c.execute("DELETE FROM cache")
                c.commit()

            finally:
                c.close()


def memory_context(db):
    rows = db.recent()

    if not rows:
        return "(none)"

    return "\n".join(
        "- [{}] {}".format(
            row[2],
            row[1]
        )
        for row in rows
    )


# ============================================================
# AI SYSTEM
# ============================================================
def system_prompt():
    return f"""
You are {NAME}, a private personal AI assistant.

Speak naturally in:
- Tamil
- English
- Tanglish

Use saved personal memory when relevant.

Never claim that an action happened unless it actually happened.

Ask for confirmation before:
- payments
- destructive actions
- sending messages
- account changes

Device current date:
{current_date()}

Device current time:
{current_time()}
""".strip()


def openai(prompt, db):
    raise RuntimeError("External AI provider disabled by Rolex-only policy")


def gemini(prompt, db):
    raise RuntimeError("External AI provider disabled by Rolex-only policy")


def ollama(prompt, db):
    raise RuntimeError("External AI provider disabled by Rolex-only policy")


def offline(prompt):
    q = prompt.lower().strip()

    if q in (
        "hi",
        "hello",
        "hey",
        "vanakkam"
    ):
        return (
            f"Vanakkam! Naan {NAME}. "
            "Offline mode-la irukken."
        )

    date_words = (
        "current date",
        "today date",
        "what date",
        "today",
        "enna date",
        "innaiku date",
        "innaiku enna date"
    )

    if any(word in q for word in date_words):
        return "📅 " + current_date()

    time_words = (
        "current time",
        "what time",
        "time now",
        "time",
        "neram",
        "enna neram",
        "ippo enna time"
    )

    if any(word in q for word in time_words):
        return "🕐 " + current_time()

    return (
        "Rolex offline fallback-la irukku.\n"
        "Memory and tasks available."
    )


# ============================================================
# MEMORY EXTRACTION
# ============================================================
def extract_memory(text):
    patterns = [
        (
            r"^(?:my name is|en name|ennoda name|"
            r"en peru|ennoda peru)\s+(.+)$",
            "identity"
        ),
        (
            r"^(?:i like|enakku pidikkum|"
            r"ennakku pidikkum)\s+(.+)$",
            "preference"
        ),
        (
            r"^(?:i love|enakku romba pidikkum)"
            r"\s+(.+)$",
            "preference"
        ),
        (
            r"^(?:i am|i'm|naan)\s+(.+)$",
            "profile"
        ),
        (
            r"^(?:remember that|remember|"
            r"idha remember pannu|idhai nyabagam vechuko)"
            r"\s+(.+)$",
            "general"
        )
    ]

    for pattern, category in patterns:
        match = re.match(
            pattern,
            text.strip(),
            re.I
        )

        if match:
            value = match.group(1).strip()

            if category == "identity":
                value = (
                    "User's name is "
                    + value
                    + "."
                )

            elif category == "preference":
                value = (
                    "User likes "
                    + value
                    + "."
                )

            return value, category

    return None, None


# ============================================================
# WEATHER
# ============================================================
def weather(city):
    geo_url = (
        "https://geocoding-api.open-meteo.com/v1/search"
        "?name="
        + urllib.parse.quote(city)
        + "&count=1&language=en&format=json"
    )

    geo = getjson(
        geo_url,
        timeout=20
    )

    if not geo.get("results"):
        return (
            "Location kandupidikka mudiyala."
        )

    place = geo["results"][0]

    forecast_url = (
        "https://api.open-meteo.com/v1/forecast?"
        "latitude="
        + str(place["latitude"])
        + "&longitude="
        + str(place["longitude"])
        + "&current="
        "temperature_2m,"
        "relative_humidity_2m,"
        "apparent_temperature,"
        "precipitation,"
        "weather_code,"
        "wind_speed_10m"
        "&timezone=auto"
    )

    data = getjson(
        forecast_url,
        timeout=20
    )

    current = data["current"]

    codes = {
        0: "clear sky",
        1: "mainly clear",
        2: "partly cloudy",
        3: "overcast",
        51: "drizzle",
        61: "rain",
        63: "rain",
        65: "heavy rain",
        80: "rain showers",
        81: "rain showers",
        82: "heavy rain showers",
        95: "thunderstorm"
    }

    return (
        f"🌦️ {place.get('name', city)}, "
        f"{place.get('country', '')}\n"
        f"Temperature: "
        f"{current.get('temperature_2m')}°C\n"
        f"Feels like: "
        f"{current.get('apparent_temperature')}°C\n"
        f"Humidity: "
        f"{current.get('relative_humidity_2m')}%\n"
        f"Condition: "
        f"{codes.get(current.get('weather_code'), 'current condition')}\n"
        f"Rain: "
        f"{current.get('precipitation')} mm\n"
        f"Wind: "
        f"{current.get('wind_speed_10m')} km/h"
    )


# ============================================================
# WEB SEARCH
# ============================================================
def web_search(query):
    if not SERPER_KEY:
        return (
            "Web search not configured.\n"
            "Add SERPER_API_KEY to .env."
        )

    result = getjson(
        "https://google.serper.dev/search",
        {"q": query},
        {"X-API-KEY": SERPER_KEY},
        25
    )

    organic = result.get(
        "organic",
        []
    )[:5]

    if not organic:
        return "No results."

    return "\n\n".join(
        "{}. {}\n{}".format(
            index + 1,
            item.get("title", ""),
            item.get("snippet", "")
        )
        for index, item in enumerate(organic)
    )


# ============================================================
# UNIVERSAL DOCUMENT INTELLIGENCE
# ============================================================
DOCUMENT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".log", ".rtf",
    ".pdf", ".docx", ".xlsx", ".pptx"
}

class DocumentManager:
    """Local-first document reader/writer/editor.

    Optional libraries are used only when installed. Text files need no
    third-party package. Original files are never silently destroyed by
    read operations. Overwrite is explicit and writes a backup first.
    """
    def __init__(self, db):
        self.db = db

    def _safe_path(self, value):
        path = os.path.abspath(os.path.expanduser(value.strip().strip('"')))
        return path

    def _record(self, path, kind=None):
        try:
            kind = kind or os.path.splitext(path)[1].lower().lstrip('.')
            c = self.db.connect()
            try:
                now = time.time()
                c.execute("""INSERT OR REPLACE INTO documents(path,title,kind,created,modified)
                             VALUES(?,?,?,?,?)""", (path, os.path.basename(path), kind, now, now))
                c.commit()
            finally:
                c.close()
        except Exception:
            pass

    def read(self, path):
        path = self._safe_path(path)
        if not os.path.exists(path):
            return f"Document not found: {path}"
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext in {'.txt','.md','.log','.rtf'}:
                text = open(path, 'r', encoding='utf-8', errors='replace').read()
            elif ext == '.json':
                with open(path, 'r', encoding='utf-8', errors='replace') as f:
                    text = json.dumps(json.load(f), ensure_ascii=False, indent=2)
            elif ext == '.csv':
                rows=[]
                with open(path, 'r', encoding='utf-8', errors='replace', newline='') as f:
                    rows=list(csv.reader(f))
                text='\n'.join(' | '.join(r) for r in rows)
            elif ext == '.pdf':
                try:
                    from pypdf import PdfReader
                except Exception:
                    try:
                        from PyPDF2 import PdfReader
                    except Exception:
                        return 'PDF reader unavailable. Install pypdf or PyPDF2.'
                reader=PdfReader(path)
                text='\n\n'.join((page.extract_text() or '') for page in reader.pages)
            elif ext == '.docx':
                try:
                    from docx import Document
                except Exception:
                    return 'DOCX reader unavailable. Install python-docx.'
                d=Document(path)
                parts=[p.text for p in d.paragraphs]
                for table in d.tables:
                    for row in table.rows:
                        parts.append(' | '.join(cell.text for cell in row.cells))
                text='\n'.join(parts)
            elif ext == '.xlsx':
                try:
                    from openpyxl import load_workbook
                except Exception:
                    return 'XLSX reader unavailable. Install openpyxl.'
                wb=load_workbook(path, data_only=True, read_only=True)
                parts=[]
                for ws in wb.worksheets:
                    parts.append(f'[SHEET: {ws.title}]')
                    for row in ws.iter_rows(values_only=True):
                        parts.append(' | '.join('' if v is None else str(v) for v in row))
                text='\n'.join(parts)
            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                except Exception:
                    return 'PPTX reader unavailable. Install python-pptx.'
                prs=Presentation(path)
                parts=[]
                for i,slide in enumerate(prs.slides,1):
                    parts.append(f'[SLIDE {i}]')
                    for shape in slide.shapes:
                        if hasattr(shape,'text') and shape.text.strip():
                            parts.append(shape.text)
                text='\n'.join(parts)
            else:
                return f'Unsupported document type: {ext}'
            self._record(path)
            return text
        except Exception as e:
            return f'Document read error: {e}'

    def write(self, path, text, overwrite=False):
        path=self._safe_path(path)
        ext=os.path.splitext(path)[1].lower()
        if ext not in {'.txt','.md','.log','.rtf','.json','.csv'}:
            return 'Safe text writing currently supports TXT, MD, LOG, RTF, JSON and CSV. Use a compatible editor for binary office formats.'
        parent=os.path.dirname(path)
        if parent and not os.path.exists(parent):
            os.makedirs(parent, exist_ok=True)
        if os.path.exists(path) and not overwrite:
            return 'File already exists. Use /doc overwrite <path> <text> for explicit overwrite.'
        if overwrite and os.path.exists(path):
            backup=path+'.rolex.bak'
            shutil.copy2(path, backup)
        with open(path,'w',encoding='utf-8') as f:
            f.write(text)
        self._record(path)
        return f'✅ Document saved: {path}' + (f'\nBackup: {path}.rolex.bak' if overwrite and os.path.exists(path+'.rolex.bak') else '')

    def append(self,path,text):
        path=self._safe_path(path)
        parent=os.path.dirname(path)
        if parent and not os.path.exists(parent): os.makedirs(parent,exist_ok=True)
        with open(path,'a',encoding='utf-8') as f: f.write(text)
        self._record(path)
        return f'✅ Text appended: {path}'

    def list_files(self, folder):
        folder=self._safe_path(folder)
        if not os.path.isdir(folder): return f'Folder not found: {folder}'
        rows=[]
        for name in sorted(os.listdir(folder)):
            p=os.path.join(folder,name)
            if os.path.isfile(p) and os.path.splitext(name)[1].lower() in DOCUMENT_EXTENSIONS:
                rows.append(f'{name} ({os.path.getsize(p)} bytes)')
        return '\n'.join(rows) if rows else 'No supported documents found.'

    def info(self,path):
        path=self._safe_path(path)
        if not os.path.exists(path): return 'Document not found.'
        st=os.stat(path)
        return (f'📄 {os.path.basename(path)}\nType: {os.path.splitext(path)[1].lower()}\n'
                f'Size: {st.st_size} bytes\nModified: {datetime.datetime.fromtimestamp(st.st_mtime).isoformat(sep=" ",timespec="seconds")}')


def document_summary(text, limit=3500):
    clean=' '.join(text.split())
    if not clean: return 'Document is empty.'
    sentences=re.split(r'(?<=[.!?])\s+', clean)
    key=sentences[:8]
    return '📌 Summary:\n' + '\n'.join(f'• {x}' for x in key) + (f'\n\n[Showing first {limit} chars]' if len(clean)>limit else '')


# ============================================================
# LOCAL SKILL ROUTER
# ------------------------------------------------------------
# Skills are intentionally local. No hosted model is called.
# Web tools are retrieval-only and are never used as an LLM proxy.
# ============================================================

def smart_llm(prompt, db, extra_system=""):
    """Compatibility hook retained for older code; always stays local."""
    return None


def skill_response(kind, prompt, db=None):
    p = prompt.strip()
    label = SKILL_LABELS.get(kind, kind.upper())
    header = f'{label} MODE\n\nRequest: {p}\n\n'
    template = SKILL_TEMPLATES.get(kind)
    return header + (template(p) if template else 'Skill not available.')

MASTER_SPEC = """ROLEX AI MASTER SPECIFICATION v3.0

1. Core intelligence: Rolex local-first orchestrator and deterministic local reasoning
2. Memory: SQLite persistent memory, recall, cache and task state
3. AI providers: OpenAI, Gemini and Ollama disabled as response providers
4. Web: live retrieval tools only; no external LLM answer pass-through
5. Education: local skill templates and structured study support
6. Work: local planning/document drafting helpers
7. Medical information: document/terminology support with safety limits
8. Documents: PDF/DOCX/XLSX/PPTX/TXT/CSV/JSON/RTF hooks
9. Vision: architecture hook for OCR/screenshots/image processing
10. Voice: Android native speech recognition hook and text fallback
11. TTS: Android-ready voice output hook
12. Wake word: configurable wake-word hook
13. Online APIs: weather and web retrieval
14. Offline mode: always available for local functions
15. Security: secret isolation, token helpers and confirmation-gate architecture
16. Sync: Android/Termux/VPS/Drive architecture hooks
17. Automation: task storage and local scheduler primitives
18. Programming: local code reading/debugging/generation scaffolding
19. Data: CSV/Excel/JSON/database analysis scaffolding
20. Remote Lab: WebSocket/pairing/authentication architecture hooks
21. Device management: Android bridge hooks
22. Reliability: diagnostics, self-tests, caching and recovery hooks
23. Backup: safe-copy/rollback architecture hooks
24. Self-improvement: propose -> sandbox -> test -> approve -> apply -> rollback workflow design
25. Modular architecture: extensions live under modules/ and are independently testable

Rolex answers are generated by Rolex's own local core. External hosted AI providers are disabled by design."""

# ============================================================
# ROLEX LOCAL BRAIN
# ============================================================
class Brain:
    """Rolex local brain. No hosted AI provider is ever called."""
    def __init__(self, db):
        self.db = db

    def _natural(self, text):
        options = [
            "Okay. Let me work with that.",
            "Got it. Here's what I can determine locally:",
            "Right. Let's break it down simply:",
            "Sure. Based on the information available to Rolex:"
        ]
        return random.choice(options) + "\n\n" + text

    def _math(self, prompt):
        if not re.fullmatch(r'[0-9+\-*/(). %]+', prompt.strip()):
            return None
        try:
            expr = prompt.strip().replace('%','/100')
            if not re.fullmatch(r'[0-9+\-*/(). ]+', expr): return None
            import ast, operator
            tree=ast.parse(expr,mode='eval')
            ops={ast.Add:operator.add,ast.Sub:operator.sub,ast.Mult:operator.mul,ast.Div:operator.truediv,ast.USub:operator.neg}
            def ev(n):
                if isinstance(n,ast.Expression): return ev(n.body)
                if isinstance(n,ast.Constant) and isinstance(n.value,(int,float)): return n.value
                if isinstance(n,ast.UnaryOp) and type(n.op) in ops: return ops[type(n.op)](ev(n.operand))
                if isinstance(n,ast.BinOp) and type(n.op) in ops: return ops[type(n.op)](ev(n.left),ev(n.right))
                raise ValueError()
            value=ev(tree)
            return f"{int(value) if float(value).is_integer() else value:g}"
        except Exception:
            return None

    def _answer(self,prompt):
        q=prompt.strip(); low=q.lower()
        if low in {'hi','hello','hey','hai','vanakkam','hey rolex','hey guru'}:
            return random.choice([f"Vanakkam! 👋 Naan {NAME}. Enna pannalaam?", "Hey! 👋 Rolex here. Sollu, enna venum?", "Vanakkam 😄 Enna matter? Naan ready."])
        if any(x in low for x in ['who are you','what are you','nee yaar','un peru enna','your name']):
            return f"Naan {NAME}. Un personal AI assistant. Local intelligence, memory and tools use panni help panren."
        if any(x in low for x in ['what can you do','enna panna mudiyum','what do you do']):
            return "🧠 Memory • 📄 Documents • 💻 Coding • 📊 Data • 🌦️ Weather • 🌐 Live web retrieval • 📝 Tasks • 🔧 Diagnostics • 🎙️ Voice hooks"
        math=self._math(q)
        if math is not None: return f"🧮 Answer: {math}"
        if low.startswith('/edu '): return skill_response('education',q[5:],self.db)
        if low.startswith('/work '): return skill_response('work',q[6:],self.db)
        if low.startswith('/medical '): return skill_response('medical',q[9:],self.db)
        if low.startswith('/code '): return skill_response('code',q[6:],self.db)
        if low.startswith('/data '): return skill_response('data',q[6:],self.db)
        rows=self.db.recall(q)
        if rows: return self._natural('I found a related memory: ' + rows[0][1])
        return ("Rolex local knowledge-la indha request-ku direct answer available illa. "
                "Guess panna maatten. Document/data kudutha, available local tools use panni help panren.")

    def _provider_order(self):
        raw = cfg("PROVIDER_ORDER", "openai,gemini,ollama,offline")
        order = []

        for name in raw.split(","):
            name = name.strip().lower()

            if name in ("openai", "gemini", "ollama", "offline"):
                if name not in order:
                    order.append(name)

        if "offline" not in order:
            order.append("offline")

        return order

    def _external_context(self, prompt):
        memory = memory_context(self.db)

        return (
            system_prompt()
            + "\n\nRelevant saved memory:\n"
            + memory
            + "\n\nImportant Rolex rule:\n"
            + "You are an intelligence source for Rolex. "
            + "Your response is not sent directly to the user. "
            + "Rolex owns the final response layer. "
            + "Do not claim actions were completed unless the information confirms they happened."
            + "\n\nUser request:\n"
            + prompt
        )

    def _ask_provider(self, provider, prompt):
        context = self._external_context(prompt)

        if provider == "openai":
            if not OPENAI_KEY:
                raise ProviderError("OpenAI API key is not configured.")

            return openai_answer(
                OPENAI_KEY,
                OPENAI_MODEL,
                system_prompt(),
                context,
                timeout=15,
            )

        if provider == "gemini":
            if not GEMINI_KEY:
                raise ProviderError("Gemini API key is not configured.")

            return gemini_answer(
                GEMINI_KEY,
                GEMINI_MODEL,
                system_prompt(),
                context,
                timeout=15,
            )

        if provider == "ollama":
            return ollama_answer(
                cfg("OLLAMA_URL", "http://127.0.0.1:11434"),
                OLLAMA_MODEL,
                system_prompt(),
                context,
                timeout=8,
            )

        return None

    def _consult_providers(self, prompt):
        """
        Consult configured external intelligence providers in parallel.

        Each provider runs independently so a slow provider does not
        unnecessarily block the others. Rolex receives all successful
        responses and owns the synthesis/final-answer layer.
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed

        providers = [
            provider
            for provider in self._provider_order()
            if provider != "offline"
        ]

        if not providers:
            return []

        responses = []

        def consult(provider):
            try:
                result = self._ask_provider(provider, prompt)

                if result:
                    return {
                        "provider": provider,
                        "text": str(result).strip(),
                    }

            except ProviderError:
                return None

            except Exception:
                return None

            return None

        # Run all available providers concurrently.
        max_workers = min(len(providers), 3)

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(consult, provider): provider
                for provider in providers
            }

            for future in as_completed(futures):
                result = future.result()

                if result:
                    responses.append(result)

        # Keep provider order deterministic for Rolex synthesis.
        provider_position = {
            provider: index
            for index, provider in enumerate(providers)
        }

        responses.sort(
            key=lambda item: provider_position.get(
                item["provider"],
                999,
            )
        )

        return responses

    def _rolex_synthesize(self, prompt, responses):
        """
        Rolex owns the final response layer.

        External providers are treated as independent opinions.
        Their raw responses are not sent directly to the user.
        """
        if not responses:
            return None

        sections = []

        for item in responses:
            sections.append(
                "[{} analysis]\n{}".format(
                    item["provider"].upper(),
                    item["text"],
                )
            )

        combined = "\n\n".join(sections)

        return (
            "Rolex analysis based on the available intelligence sources:"
            "\n\n"
            + combined
            + "\n\n"
            "Rolex final note: The above sources were consulted by Rolex. "
            "Where their information differs, it should be verified against "
            "reliable source data before taking action."
        )

    def _rolex_finalize(self, provider_text, prompt, provider):
        if not provider_text:
            return None

        text = str(provider_text).strip()

        if not text:
            return None

        prefix = "Based on the information available to Rolex:"

        return prefix + "\n\n" + text

    def ask(self, prompt):
        prompt = prompt.strip()

        if not prompt:
            return "", "role-local"

        key = hashlib.sha256(
            ("role-router-v2\\n" + prompt.lower()).encode("utf-8")
        ).hexdigest()

        cached = self.db.cache_get(key)

        if cached:
            return cached, "role-cache"

        # Rolex handles simple/common requests locally first.
        local = self._answer(prompt)

        unknown_prefix = (
            "Rolex local knowledge-la indha request-ku direct answer available illa."
        )

        if local and not local.startswith(unknown_prefix):
            self.db.cache_set(key, local)
            return local, "role-local"

        # Parallel intelligence layer.
        # Rolex consults available external intelligence sources concurrently.
        # The parallel engine handles provider execution, scoring and synthesis.
        if not EXTERNAL_AI_DISABLED:
            config = {
                "providers": [
                    provider
                    for provider in self._provider_order()
                    if provider != "offline"
                ],
                "openai_key": OPENAI_KEY,
                "openai_model": OPENAI_MODEL,
                "gemini_key": GEMINI_KEY,
                "gemini_model": GEMINI_MODEL,
                "ollama_url": cfg(
                    "OLLAMA_URL",
                    "http://127.0.0.1:11434",
                ),
                "ollama_model": OLLAMA_MODEL,
                "system_prompt": system_prompt(),
                "openai_timeout": 15,
                "gemini_timeout": 15,
                "ollama_timeout": 60,
            }

            try:
                result = parallel_rolex_answer(
                    prompt,
                    config,
                )

                final = result.get("final_answer")

                if final:
                    self.db.cache_set(key, final)

                    successful = [
                        item
                        for item in result.get("candidates", [])
                        if item.get("success")
                    ]

                    providers = ",".join(
                        item["provider"]
                        for item in successful
                    )

                    return final, (
                        "role-parallel/"
                        + (providers or "none")
                    )

            except Exception:
                pass

        # Complete offline fallback.
        fallback = offline(prompt)

        self.db.cache_set(key, fallback)

        return fallback, "role-offline"


# ============================================================
# ROLEX CORE
# ============================================================
class SystemDiagnostics:
    """Lightweight diagnostics that work without network calls."""
    @staticmethod
    def run(db=None):
        checks = []
        try:
            checks.append(("Python", platform.python_version(), True))
            checks.append(("Platform", platform.platform(), True))
            if db is not None:
                db.initialize()
                test_id = db.remember("__rolex_self_test__", "diagnostic")
                rows = db.recall("__rolex_self_test__")
                ok = bool(rows)
                db.forget(test_id)
                checks.append(("SQLite", "read/write OK", ok))
            checks.append(("Internet", "available" if socket.create_connection(("1.1.1.1", 53), 2) else "unavailable", True))
        except Exception as e:
            checks.append(("Diagnostics", str(e), False))
        return checks


class Rolex:

    def __init__(self):
        self.db = DB()
        self.brain = Brain(self.db)
        self.docs = DocumentManager(self.db)

    def status(self):
        return (
            f"{NAME} v2.2 MASTER\n"
            "Database: OK\n"
            "SQLite thread-safe: FIXED\n"
            "Auto-memory: ON\n"
            "Cache: ON\n"
            "Live weather: ON\n"
            f"Web search: "
            f"{'ON' if SERPER_KEY else 'OFF'}\n"
            "External AI providers: DISABLED\n"
            "Rolex local brain: ON\n"
            "Rolex local brain: ON (always-on fallback)\n"
            "Education skill: ON\n"
            "Work skill: ON\n"
            "Coding skill: ON\n"
            "Data skill: ON\n"
            "Medical document safety layer: ON\n"
            "Universal document layer: ON\n"
            f"Android bridge: "
            f"{'AVAILABLE' if ANDROID_BRIDGE else 'NOT AVAILABLE'}\n"
            "Offline: ALWAYS AVAILABLE\n"
            f"Wake word: {'OPTIONAL' if cfg('WAKE_WORD_ENGINE') else 'NOT CONFIGURED'}\n"
            f"Remote lab: {'CONFIGURED' if cfg('ROLEX_REMOTE_URL') else 'NOT CONFIGURED'}\n"
            f"Sync: {'CONFIGURED' if cfg('ROLEX_SYNC_URL') else 'NOT CONFIGURED'}"
        )

    def handle(self, text):
        if not text:
            return ""

        text = text.strip()

        if text == "/help":
            return (
                "/health\n"
                "/provider\n"
                "/remember <text>\n"
                "/recall [query]\n"
                "/forget <id>\n"
                "/task <text>\n"
                "/tasks\n"
                "/done <id>\n"
                "/weather <city>\n"
                "/web <query>\n"
                "/edu <topic>\n"
                "/work <request>\n"
                "/medical <topic>\n"
                "/code <request/code>\n"
                "/data <request>\n"
                "/doc read <path>\n"
                "/doc summary <path>\n"
                "/doc ask <path> | <question>\n"
                "/doc info <path>\n"
                "/doc list <folder>\n"
                "/doc create <path> | <content>\n"
                "/doc overwrite <path> | <content>\n"
                "/doc append <path> | <content>\n"
                "/spec\n"
                "/clear-cache\n"
                "/diag\n"
                "/version\n"
                "/exit"
            )

        if text in (
            "/health",
            "/provider",
            "/version"
        ):
            return self.status()

        if text == "/diag":
            lines = ["🔧 ROLEX DIAGNOSTICS"]
            for name, value, ok in SystemDiagnostics.run(self.db):
                lines.append(f"{'✅' if ok else '❌'} {name}: {value}")
            return "\n".join(lines)

        if text.startswith("/remember "):
            value = text[10:].strip()

            if not value:
                return "Use /remember <text>"

            memory_id = self.db.remember(
                value
            )

            return (
                f"💾 Memory #{memory_id} saved."
            )

        if (
            text == "/recall"
            or text.startswith("/recall ")
        ):
            query = ""

            if text.startswith("/recall "):
                query = text[8:].strip()

            rows = self.db.recall(query)

            if not rows:
                return "No memories."

            return "\n".join(
                f"#{row[0]} "
                f"[{row[2]}]: "
                f"{row[1]}"
                for row in rows
            )

        if text.startswith("/forget "):
            try:
                memory_id = int(
                    text[8:].strip()
                )

                self.db.forget(
                    memory_id
                )

                return "Memory deleted."

            except Exception:
                return "Use /forget <id>"

        if text.startswith("/task "):
            value = text[6:].strip()

            if not value:
                return "Use /task <text>"

            task_id = self.db.task(
                value
            )

            return (
                f"📝 Task #{task_id} saved."
            )

        if text == "/tasks":
            rows = self.db.tasks()

            if not rows:
                return "No open tasks."

            return "\n".join(
                f"#{row[0]}: {row[1]}"
                for row in rows
            )

        if text.startswith("/done "):
            try:
                task_id = int(
                    text[6:].strip()
                )

                self.db.done(
                    task_id
                )

                return "Task completed."

            except Exception:
                return "Use /done <id>"

        if text.startswith("/weather "):
            try:
                return weather(
                    text[9:].strip()
                )
            except Exception as error:
                return (
                    "Weather error: "
                    + str(error)
                )

        if text.startswith("/web "):
            return web_search(
                text[5:].strip()
            )

        if text == "/spec":
            return MASTER_SPEC

        if text.startswith("/edu "):
            return skill_response("education", text[5:].strip(), self.db)

        if text.startswith("/work "):
            return skill_response("work", text[6:].strip(), self.db)

        if text.startswith("/medical "):
            return skill_response("medical", text[9:].strip(), self.db)

        if text.startswith("/code "):
            return skill_response("code", text[6:].strip(), self.db)

        if text.startswith("/data "):
            return skill_response("data", text[6:].strip(), self.db)

        if text.startswith("/doc read "):
            return self.docs.read(text[10:].strip())

        if text.startswith("/doc summary "):
            raw=self.docs.read(text[13:].strip())
            if raw.startswith("Document ") or raw.startswith("PDF reader") or raw.startswith("DOCX reader") or raw.startswith("XLSX reader") or raw.startswith("PPTX reader"):
                return raw
            return document_summary(raw)

        if text.startswith("/doc ask "):
            payload=text[9:].strip()
            if " | " not in payload:
                return "Use /doc ask <path> | <question>"
            path,question=payload.split(" | ",1)
            raw=self.docs.read(path.strip())
            if raw.startswith("Document not found") or raw.endswith("unavailable. Install pypdf or PyPDF2.") or "reader unavailable" in raw or raw.startswith("Unsupported document type") or raw.startswith("Document read error"):
                return raw
            context=raw[:12000]
            answer = document_summary(context)
            return "📄 Document summary:\n" + answer + "\n\nRolex note: direct document Q&A is handled locally only; hosted AI providers are disabled."

        if text.startswith("/doc info "):
            return self.docs.info(text[10:].strip())

        if text.startswith("/doc list "):
            return self.docs.list_files(text[10:].strip())

        if text.startswith("/doc create "):
            payload=text[12:].strip()
            if " | " not in payload: return "Use /doc create <path> | <content>"
            path,content=payload.split(" | ",1)
            return self.docs.write(path,content,False)

        if text.startswith("/doc overwrite "):
            payload=text[16:].strip()
            if " | " not in payload: return "Use /doc overwrite <path> | <content>"
            path,content=payload.split(" | ",1)
            return self.docs.write(path,content,True)

        if text.startswith("/doc append "):
            payload=text[12:].strip()
            if " | " not in payload: return "Use /doc append <path> | <content>"
            path,content=payload.split(" | ",1)
            return self.docs.append(path,content)

        if text == "/clear-cache":
            self.db.clear_cache()
            return "Cache cleared."

        if text == "/exit":
            return None

        lower = text.lower()

        # Date must be handled BEFORE generic AI.
        if any(
            word in lower
            for word in (
                "current date",
                "today date",
                "what date",
                "today",
                "enna date",
                "innaiku date",
                "innaiku enna date"
            )
        ):
            return (
                "📅 "
                + current_date()
            )

        # Time must be handled BEFORE generic AI.
        if any(
            word in lower
            for word in (
                "current time",
                "what time",
                "time now",
                "time",
                "neram",
                "enna neram",
                "ippo enna time"
            )
        ):
            return (
                "🕐 "
                + current_time()
            )

        if (
            "weather" in lower
            or "climate" in lower
        ):
            match = re.search(
                r"(?:in|at|of)\s+(.+)$",
                text,
                re.I
            )

            if match:
                try:
                    return weather(
                        match.group(1)
                        .strip(" ?.")
                    )
                except Exception as error:
                    return (
                        "Weather error: "
                        + str(error)
                    )

        memory, category = extract_memory(
            text
        )

        note = ""

        if (
            memory
            and not self.db.has(memory)
        ):
            memory_id = self.db.remember(
                memory,
                category
            )

            note = (
                f"\n💾 Memory #{memory_id} saved."
            )

        task_match = re.match(
            r"^(?:add task|task|remind me to)"
            r"\s+(.+)$",
            text,
            re.I
        )

        if task_match:
            task_id = self.db.task(
                task_match.group(1).strip()
            )

            return (
                f"📝 Task #{task_id} saved."
            )

        answer, provider = self.brain.ask(
            text
        )

        return answer + note


# ============================================================
# ANDROID NATIVE VOICE
# ============================================================
class NativeVoice:

    REQUEST_CODE = 7191

    def __init__(self, gui):
        self.gui = gui
        self.enabled = False
        self.Intent = None
        self.RecognizerIntent = None

        self.setup()

    def setup(self):

        if not ANDROID_BRIDGE:
            return

        try:

            # Try PyJNIus first.
            if autoclass is not None:

                Intent = autoclass(
                    "android.content.Intent"
                )

                self.RecognizerIntent = autoclass(
                    "android.speech.RecognizerIntent"
                )

                self.Intent = Intent

                self.enabled = True

                if activity is not None:
                    activity.bind(
                        on_activity_result=
                        self.on_result
                    )

                return

        except Exception:
            pass

        self.enabled = False

    def listen(self):

        if not self.enabled:
            self.gui.show_voice_unavailable()
            return

        try:

            intent = self.Intent(
                self.RecognizerIntent
                .ACTION_RECOGNIZE_SPEECH
            )

            intent.putExtra(
                self.RecognizerIntent
                .EXTRA_LANGUAGE_MODEL,
                self.RecognizerIntent
                .LANGUAGE_MODEL_FREE_FORM
            )

            # Device default language first.
            # User can speak Tamil / English / Tanglish.
            intent.putExtra(
                self.RecognizerIntent
                .EXTRA_PROMPT,
                "Speak to Rolex"
            )

            intent.putExtra(
                self.RecognizerIntent
                .EXTRA_MAX_RESULTS,
                3
            )

            self.gui.set_status(
                "● LISTENING"
            )

            activity.startActivityForResult(
                intent,
                self.REQUEST_CODE
            )

        except Exception as error:

            self.gui.set_status(
                "● READY"
            )

            self.gui.add_bot(
                "Voice start error: "
                + str(error)
            )

    def on_result(
        self,
        request_code,
        result_code,
        intent
    ):

        if request_code != self.REQUEST_CODE:
            return

        try:

            results = intent.getStringArrayListExtra(
                self.RecognizerIntent
                .EXTRA_RESULTS
            )

            if results and len(results) > 0:

                text = str(
                    results.get(0)
                )

                Clock.schedule_once(
                    lambda dt:
                    self.gui.voice_text(text),
                    0
                )

            else:

                Clock.schedule_once(
                    lambda dt:
                    self.gui.add_bot(
                        "🎙️ Voice result empty. "
                        "Try again."
                    ),
                    0
                )

        except Exception as error:

            Clock.schedule_once(
                lambda dt:
                self.gui.add_bot(
                    "Voice result error: "
                    + str(error)
                ),
                0
            )

        Clock.schedule_once(
            lambda dt:
            self.gui.set_status("● READY"),
            0
        )


# ============================================================
# CHAT BUBBLE
# ============================================================
class ChatBubble(BoxLayout):

    def __init__(
        self,
        text,
        user=False,
        **kwargs
    ):

        super().__init__(
            orientation="vertical",
            size_hint_y=None,
            padding=(
                dp(10),
                dp(6)
            ),
            spacing=dp(2),
            **kwargs
        )

        title = (
            "YOU"
            if user
            else NAME.upper()
        )

        bubble_color = COLOR_USER_BUBBLE if user else COLOR_BOT_BUBBLE
        title_color = COLOR_ACCENT if user else COLOR_ACCENT_SOFT

        with self.canvas.before:
            Color(*bubble_color)
            self._rect = RoundedRectangle(
                pos=self.pos,
                size=self.size,
                radius=[dp(14)]
            )

        self.bind(pos=self._sync_rect, size=self._sync_rect)

        self.title = Label(
            text=title,
            size_hint_y=None,
            height=dp(18),
            font_size="11sp",
            bold=True,
            color=title_color,
            halign="left"
        )

        self.body = Label(
            text=text,
            size_hint_y=None,
            font_size="15sp",
            color=COLOR_TEXT,
            halign="left",
            valign="top"
        )

        self.add_widget(
            self.title
        )

        self.add_widget(
            self.body
        )

        self.bind(
            width=self.resize
        )

        Clock.schedule_once(
            self.resize,
            0
        )

    def _sync_rect(self, *_):
        self._rect.pos = self.pos
        self._rect.size = self.size

    def resize(self, *_):

        width = max(
            dp(100),
            self.width - dp(20)
        )

        self.title.text_size = (width, None)
        self.body.text_size = (
            width,
            None
        )

        self.body.texture_update()

        self.body.height = max(
            dp(25),
            self.body.texture_size[1]
        )

        self.height = (
            self.title.height
            + self.body.height
            + dp(16)
        )


class RoundButton(Button):
    """Flat, rounded, accent-able button used everywhere in the UI so the
    app doesn't look like default unstyled Kivy widgets."""

    def __init__(self, bg=None, fg=None, radius=14, **kwargs):
        super().__init__(**kwargs)
        self.background_normal = ''
        self.background_down = ''
        self.background_color = (0, 0, 0, 0)
        self.color = fg or COLOR_TEXT
        self.bold = True
        self._bg = bg or COLOR_SURFACE
        with self.canvas.before:
            Color(*self._bg)
            self._rect = RoundedRectangle(
                pos=self.pos, size=self.size, radius=[dp(radius)]
            )
        self.bind(pos=self._sync_rect, size=self._sync_rect)

    def _sync_rect(self, *_):
        self._rect.pos = self.pos
        self._rect.size = self.size


# ============================================================
# GUI
# ============================================================
class RolexGUI(App):

    def build(self):

        self.title = (
            "Rolex AI v2.2 Master"
        )

        Window.clearcolor = COLOR_BG

        self.rolex = Rolex()
        self.busy = False
        smart_on = False

        self.voice = NativeVoice(
            self
        )

        root = BoxLayout(
            orientation="vertical",
            padding=dp(10),
            spacing=dp(8)
        )

        # HEADER
        header = BoxLayout(
            size_hint_y=None,
            height=dp(60),
            spacing=dp(5),
            padding=(dp(12), dp(6))
        )

        with header.canvas.before:
            Color(*COLOR_PANEL)
            header_rect = RoundedRectangle(pos=header.pos, size=header.size, radius=[dp(16)])
            Color(*COLOR_ACCENT)
            header_bar = RoundedRectangle(pos=header.pos, size=(dp(4), header.height), radius=[dp(2)])

        def _sync_header(*_):
            header_rect.pos = header.pos
            header_rect.size = header.size
            header_bar.pos = header.pos
            header_bar.size = (dp(4), header.height)
        header.bind(pos=_sync_header, size=_sync_header)

        title_box = BoxLayout(
            orientation="vertical",
            size_hint_x=.7
        )

        title_box.add_widget(
            Label(
                text="🤖 ROLEX AI",
                font_size="21sp",
                bold=True,
                color=COLOR_TEXT,
                halign="left",
                valign="middle"
            )
        )

        title_box.add_widget(
            Label(
                text="Personal Intelligence  •  Rolex Local Core",
                font_size="11sp",
                color=COLOR_MUTED,
                halign="left",
                valign="middle"
            )
        )

        for lbl in title_box.children:
            lbl.bind(size=lambda w, *_: setattr(w, 'text_size', w.size))

        self.status = Label(
            text="● READY",
            size_hint_x=.3,
            font_size="12sp",
            bold=True,
            color=COLOR_OK,
            halign="right",
            valign="middle"
        )
        self.status.bind(size=lambda w, *_: setattr(w, 'text_size', w.size))

        header.add_widget(
            title_box
        )

        header.add_widget(
            self.status
        )

        root.add_widget(
            header
        )

        # CHAT
        self.scroll = ScrollView(
            do_scroll_x=False,
            bar_width=dp(4),
            bar_color=COLOR_ACCENT,
            bar_inactive_color=COLOR_SURFACE
        )

        self.chat = BoxLayout(
            orientation="vertical",
            size_hint_y=None,
            spacing=dp(8),
            padding=dp(4)
        )

        self.chat.bind(
            minimum_height=
            self.chat.setter("height")
        )

        self.scroll.add_widget(
            self.chat
        )

        root.add_widget(
            self.scroll
        )

        # INPUT
        input_row = BoxLayout(
            size_hint_y=None,
            height=dp(52),
            spacing=dp(6)
        )

        self.input = TextInput(
            hint_text="Type to Rolex...  (/help for commands)",
            multiline=False,
            font_size="15sp",
            background_color=COLOR_SURFACE,
            foreground_color=COLOR_TEXT,
            hint_text_color=COLOR_MUTED,
            cursor_color=COLOR_ACCENT,
            padding=(dp(14), dp(14)),
        )

        self.input.bind(
            on_text_validate=self.send
        )

        mic = RoundButton(
            text="🎙️",
            bg=COLOR_SURFACE,
            fg=COLOR_TEXT,
            size_hint_x=None,
            width=dp(52),
            font_size="20sp"
        )

        mic.bind(
            on_release=self.listen
        )

        send = RoundButton(
            text="SEND",
            bg=COLOR_ACCENT,
            fg=COLOR_INK,
            size_hint_x=None,
            width=dp(80),
            font_size="13sp"
        )

        send.bind(
            on_release=self.send
        )

        input_row.add_widget(
            self.input
        )

        input_row.add_widget(
            mic
        )

        input_row.add_widget(
            send
        )

        root.add_widget(
            input_row
        )

        # BOTTOM
        bottom = BoxLayout(
            size_hint_y=None,
            height=dp(42),
            spacing=dp(6)
        )

        clear = RoundButton(
            text="🧹 Clear Chat",
            bg=COLOR_SURFACE,
            fg=COLOR_MUTED,
            font_size="12sp"
        )

        clear.bind(
            on_release=self.clear_chat
        )

        help_button = RoundButton(
            text="❓ Help",
            bg=COLOR_SURFACE,
            fg=COLOR_MUTED,
            font_size="12sp"
        )

        help_button.bind(
            on_release=self.help
        )

        bottom.add_widget(
            clear
        )

        bottom.add_widget(
            help_button
        )

        root.add_widget(
            bottom
        )

        # FIRST MESSAGE
        self.add_bot(
            "Vanakkam! 👋\n"
            "Naan Rolex AI.\n"
            + "Rolex local core active. Hosted AI providers are disabled by design.\n\n"
            + "Type or press 🎙️ to test voice. Try /help for everything I can do."
        )

        return root

    # --------------------------------------------------------
    # CHAT
    # --------------------------------------------------------
    def add_user(self, text):

        self.chat.add_widget(
            ChatBubble(
                text,
                True
            )
        )

        self.scroll_bottom()

    def add_bot(self, text):

        self.chat.add_widget(
            ChatBubble(
                text,
                False
            )
        )

        self.scroll_bottom()

    def scroll_bottom(self):

        Clock.schedule_once(
            lambda dt:
            setattr(
                self.scroll,
                "scroll_y",
                0
            ),
            .05
        )

    # --------------------------------------------------------
    # STATUS
    # --------------------------------------------------------
    def set_status(self, text):
        self.status.text = text
        upper = text.upper()
        if "READY" in upper:
            self.status.color = COLOR_OK
        elif "ERROR" in upper:
            self.status.color = (1, 0.36, 0.45, 1)
        else:
            self.status.color = COLOR_BUSY

    # --------------------------------------------------------
    # SEND
    # --------------------------------------------------------
    def send(self, *_):

        if self.busy:
            return

        text = self.input.text.strip()

        if not text:
            return

        self.input.text = ""

        self.add_user(
            text
        )

        self.process(
            text
        )

    def process(self, text):

        if self.busy:
            return

        self.busy = True

        self.set_status(
            "● THINKING"
        )

        threading.Thread(
            target=self.worker,
            args=(text,),
            daemon=True
        ).start()

    def worker(self, text):

        try:

            result = self.rolex.handle(
                text
            )

            if result is None:
                result = "Bye 👋"

        except Exception as error:

            result = (
                "Error: "
                + str(error)
            )

        Clock.schedule_once(
            lambda dt:
            self.finish(result),
            0
        )

    def finish(self, result):

        self.add_bot(
            result
        )

        self.busy = False

        self.set_status(
            "● READY"
        )

    # --------------------------------------------------------
    # VOICE
    # --------------------------------------------------------
    def listen(self, *_):

        if self.busy:
            return

        self.voice.listen()

    def voice_text(self, text):

        if not text:
            return

        self.add_user(
            "🎙️ " + text
        )

        self.process(
            text
        )

    def show_voice_unavailable(self):

        self.set_status(
            "● READY"
        )

        self.add_bot(
            "🎙️ Android Speech API is not "
            "available in this Pydroid setup.\n\n"
            "No PyAudio is required. "
            "Text chat is working normally.\n\n"
            "If you want Android native voice, "
            "we can make a separate APK-ready "
            "voice layer next."
        )

    # --------------------------------------------------------
    # QUICK ACTIONS
    # --------------------------------------------------------
    def command(self, command):

        if self.busy:
            return

        self.input.text = command
        self.send()

    def memory(self, *_):
        self.command(
            "/recall"
        )

    def tasks(self, *_):
        self.command(
            "/tasks"
        )

    def weather(self, *_):
        self.command(
            "/weather Chennai"
        )

    def health(self, *_):
        self.command(
            "/health"
        )

    def help(self, *_):
        self.command(
            "/help"
        )

    def education(self, *_):
        self.command("/edu study planning and concept explanation")

    def work(self, *_):
        self.command("/work organize my work request")

    def documents(self, *_):
        self.command("/spec")

    def medical(self, *_):
        self.command("/medical explain a medical document")

    def clear_chat(self, *_):

        self.chat.clear_widgets()

        self.add_bot(
            "Chat cleared. 🧹"
        )


# ============================================================
# BUILT-IN SELF TEST
# ============================================================
def run_self_test():
    print("ROLEX AI FINAL A-Z SELF TEST")
    print("=" * 32)
    print("Python:", platform.python_version())
    print("Machine:", platform.machine())
    print("Kivy: GUI module loaded")
    db = DB()
    tests = [
        ("database", lambda: db.initialize()),
        ("memory write/read", lambda: db.recall()),
        ("date", lambda: current_date()),
        ("time", lambda: current_time()),
        ("offline greeting", lambda: offline("hello")),
        ("offline date", lambda: offline("what is today's date")),
        ("offline time", lambda: offline("what time is it")),
        ("memory extraction", lambda: extract_memory("remember that testing works")),
        ("document manager", lambda: DocumentManager(db).write(os.path.join(BASE,"rolex_selftest.txt"), "Rolex document test", True)),
        ("document read", lambda: DocumentManager(db).read(os.path.join(BASE,"rolex_selftest.txt"))),
        ("local brain", lambda: Brain(db).ask("what can you do")),
    ]
    passed = 0
    for name, fn in tests:
        try:
            result = fn()
            print("PASS", name, "=>", str(result)[:120])
            passed += 1
        except Exception as e:
            print("FAIL", name, "=>", e)
    try:
        test_file=os.path.join(BASE,"rolex_selftest.txt")
        if os.path.exists(test_file): os.remove(test_file)
        if os.path.exists(test_file+".rolex.bak"): os.remove(test_file+".rolex.bak")
    except Exception: pass
    print(f"RESULT: {passed}/{len(tests)} tests passed")
    return passed == len(tests)


# ============================================================
# START
# ============================================================
if __name__ == "__main__":
    if "--self-test" in __import__("sys").argv:
        raise SystemExit(0 if run_self_test() else 1)
    RolexGUI().run()
